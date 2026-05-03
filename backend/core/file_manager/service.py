#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
文件管理服务
"""
import base64
import mimetypes
import os
from typing import Optional, List, Tuple

from sqlalchemy import select, func, update, case, literal_column, or_
from sqlalchemy.ext.asyncio import AsyncSession

from app.base_service import BaseService
from utils.context import get_current_user_id_from_context
from core.file_manager.model import FileManager
from core.file_manager.schema import FileManagerCreate, FileManagerUpdate
from core.file_manager.storage_backends import get_storage_backend, LocalStorageBackend, MinioStorageBackend


class FileManagerService(BaseService[FileManager, FileManagerCreate, FileManagerUpdate]):
    """文件管理服务"""
    
    model = FileManager

    @classmethod
    async def get_list(
        cls,
        db: AsyncSession,
        page: int = 1,
        page_size: int = 20,
        parent_id: Optional[str] = None,
        name: Optional[str] = None,
        type: Optional[str] = None,
        storage_type: Optional[str] = None,
        file_ext: Optional[str] = None,
        is_public: Optional[bool] = None,
        creator_id: Optional[str] = None,
        is_superuser: bool = False,
    ) -> Tuple[List[FileManager], int]:
        """获取文件列表
        
        权限规则：
        - 超管(is_superuser=True)：看到所有文件
        - 普通用户：只看到自己创建的 + 公共的 + 系统文件夹
        """
        # 构建查询条件
        conditions = [cls.model.is_deleted == False]  # noqa: E712
        
        # 父文件夹过滤
        if parent_id is None:
            conditions.append(cls.model.parent_id == None)  # noqa: E711
        else:
            conditions.append(cls.model.parent_id == parent_id)
        
        if name:
            conditions.append(cls.model.name.ilike(f"%{name}%"))
        if type:
            conditions.append(cls.model.type == type)
        if storage_type:
            conditions.append(cls.model.storage_type == storage_type)
        if file_ext:
            conditions.append(cls.model.file_ext == file_ext)
        if is_public is not None:
            conditions.append(cls.model.is_public == is_public)
        
        # 权限过滤：普通用户只看自己的 + 公共的 + 系统文件夹
        if not is_superuser and creator_id:
            conditions.append(
                or_(
                    cls.model.sys_creator_id == creator_id,
                    cls.model.is_public == True,  # noqa: E712
                    cls.model.is_system == True,  # noqa: E712
                )
            )
        
        # 查询总数
        count_query = select(func.count(cls.model.id)).where(*conditions)
        total_result = await db.execute(count_query)
        total = total_result.scalar() or 0
        
        # 查询数据（文件夹排在前面，新建的排最前）
        offset = (page - 1) * page_size
        type_order = case(
            (cls.model.type == 'folder', 0),
            else_=1
        )
        query = (
            select(cls.model)
            .where(*conditions)
            .order_by(type_order, cls.model.sys_create_datetime.desc())
            .offset(offset)
            .limit(page_size)
        )
        result = await db.execute(query)
        items = result.scalars().all()
        
        return items, total

    @classmethod
    async def get_folder_tree(
        cls,
        db: AsyncSession,
        creator_id: Optional[str] = None,
        is_superuser: bool = False,
    ) -> List[FileManager]:
        """获取文件夹树结构
        
        权限规则：
        - 超管：看到所有文件夹
        - 普通用户：只看到自己创建的 + 公共的 + 系统文件夹
        """
        conditions = [
            cls.model.type == 'folder',
            cls.model.is_deleted == False,  # noqa: E712
        ]
        
        if not is_superuser and creator_id:
            conditions.append(
                or_(
                    cls.model.sys_creator_id == creator_id,
                    cls.model.is_public == True,  # noqa: E712
                    cls.model.is_system == True,  # noqa: E712
                )
            )
        
        query = (
            select(cls.model)
            .where(*conditions)
            .order_by(cls.model.name)
        )
        result = await db.execute(query)
        return result.scalars().all()

    IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.svg', '.ico', '.tiff', '.heic', '.heif'}

    @classmethod
    async def get_recent_images(
        cls,
        db: AsyncSession,
        creator_id: str,
        limit: int = 20,
    ) -> List[FileManager]:
        """获取当前用户最近上传的图片文件（跨所有文件夹，按时间倒序）"""
        query = (
            select(cls.model)
            .where(
                cls.model.type == 'file',
                cls.model.is_deleted == False,  # noqa: E712
                cls.model.file_ext.in_(cls.IMAGE_EXTENSIONS),
                cls.model.sys_creator_id == creator_id,
            )
            .order_by(cls.model.sys_create_datetime.desc())
            .limit(limit)
        )
        result = await db.execute(query)
        return list(result.scalars().all())

    @classmethod
    async def get_recent_files(
        cls,
        db: AsyncSession,
        creator_id: str,
        limit: int = 20,
    ) -> List[FileManager]:
        """获取当前用户最近上传的文件（跨所有文件夹，按时间倒序）"""
        query = (
            select(cls.model)
            .where(
                cls.model.type == 'file',
                cls.model.is_deleted == False,  # noqa: E712
                cls.model.sys_creator_id == creator_id,
            )
            .order_by(cls.model.sys_create_datetime.desc())
            .limit(limit)
        )
        result = await db.execute(query)
        return list(result.scalars().all())

    @classmethod
    async def create_folder(
        cls,
        db: AsyncSession,
        name: str,
        parent_id: Optional[str] = None,
        creator_id: Optional[str] = None,
    ) -> Optional[FileManager]:
        """创建文件夹"""
        if not creator_id:
            creator_id = get_current_user_id_from_context()
        # 获取父文件夹路径
        parent_path = ''
        if parent_id:
            parent = await cls.get_by_id(db, parent_id)
            if parent and parent.type == 'folder':
                parent_path = parent.path
        
        # 构建文件夹路径
        folder_path = os.path.join(parent_path, name).replace('\\', '/') if parent_path else name
        
        # 检查同名文件夹
        existing = await db.execute(
            select(cls.model).where(
                cls.model.parent_id == parent_id,
                cls.model.name == name,
                cls.model.type == 'folder',
                cls.model.is_deleted == False  # noqa: E712
            )
        )
        if existing.scalar_one_or_none():
            return None  # 同名文件夹已存在
        
        # 创建文件夹
        folder = FileManager(
            name=name,
            type='folder',
            parent_id=parent_id,
            path=folder_path,
            storage_path='',
            sys_creator_id=creator_id,
        )
        db.add(folder)
        await db.commit()
        await db.refresh(folder)
        return folder

    # 来源模块 → 系统文件夹名称映射
    SOURCE_LABELS = {
        'announcement': '公告附件',
        'workflow': '工作流附件',
        'chat': '聊天文件',
        'form': '表单附件',
        'avatar': '头像',
        'ai_platform': 'AI平台文件',
        'screen_design': '大屏设计素材',
    }

    SYSTEM_ROOT_NAME = 'SystemFile'

    @classmethod
    async def get_or_create_source_folder(
        cls,
        db: AsyncSession,
        source: str,
    ) -> str:
        """根据 source 自动获取或创建 SystemFile/模块文件夹/年-月 三级文件夹，返回日期子文件夹 ID"""
        from datetime import datetime
        folder_name = cls.SOURCE_LABELS.get(source, source)

        # 1. 获取或创建根级 SystemFile 文件夹
        result = await db.execute(
            select(cls.model).where(
                cls.model.parent_id == None,  # noqa: E711
                cls.model.name == cls.SYSTEM_ROOT_NAME,
                cls.model.type == 'folder',
                cls.model.is_system == True,  # noqa: E712
                cls.model.is_deleted == False,  # noqa: E712
            )
        )
        system_root = result.scalar_one_or_none()
        if not system_root:
            system_root = FileManager(
                name=cls.SYSTEM_ROOT_NAME,
                type='folder',
                parent_id=None,
                path=cls.SYSTEM_ROOT_NAME,
                storage_path='',
                is_system=True,
                source='system',
            )
            db.add(system_root)
            await db.flush()

        # 2. 获取或创建模块级文件夹（SystemFile 下）
        module_path = f"{cls.SYSTEM_ROOT_NAME}/{folder_name}"
        result = await db.execute(
            select(cls.model).where(
                cls.model.parent_id == system_root.id,
                cls.model.name == folder_name,
                cls.model.type == 'folder',
                cls.model.is_system == True,  # noqa: E712
                cls.model.is_deleted == False,  # noqa: E712
            )
        )
        module_folder = result.scalar_one_or_none()
        if not module_folder:
            module_folder = FileManager(
                name=folder_name,
                type='folder',
                parent_id=system_root.id,
                path=module_path,
                storage_path='',
                is_system=True,
                source=source,
            )
            db.add(module_folder)
            await db.flush()

        # 3. 获取或创建日期子文件夹（年-月）
        date_name = datetime.now().strftime('%Y-%m')
        date_path = f"{module_path}/{date_name}"

        result = await db.execute(
            select(cls.model).where(
                cls.model.parent_id == module_folder.id,
                cls.model.name == date_name,
                cls.model.type == 'folder',
                cls.model.is_deleted == False,  # noqa: E712
            )
        )
        date_folder = result.scalar_one_or_none()
        if not date_folder:
            date_folder = FileManager(
                name=date_name,
                type='folder',
                parent_id=module_folder.id,
                path=date_path,
                storage_path='',
                is_system=True,
                source=source,
            )
            db.add(date_folder)
            await db.flush()

        return date_folder.id

    @classmethod
    async def upload_file(
        cls,
        db: AsyncSession,
        file_content: bytes,
        filename: str,
        file_size: int,
        parent_id: Optional[str] = None,
        is_public: bool = False,
        creator_id: Optional[str] = None,
        source: Optional[str] = None,
    ) -> FileManager:
        """上传文件"""
        if not creator_id:
            creator_id = get_current_user_id_from_context()
        # 如果有 source 且没有指定 parent_id，自动归入系统文件夹
        if source and not parent_id:
            parent_id = await cls.get_or_create_source_folder(db, source)

        # 获取父文件夹路径
        folder_path = ''
        if parent_id:
            parent = await cls.get_by_id(db, parent_id)
            if parent and parent.type == 'folder':
                folder_path = parent.path
        
        # 获取存储后端
        storage = get_storage_backend()
        
        # 计算文件信息
        file_ext = os.path.splitext(filename)[1].lower()
        mime_type = mimetypes.guess_type(filename)[0] or 'application/octet-stream'
        
        # 创建文件对象用于保存
        import io
        file_obj = io.BytesIO(file_content)
        
        # 计算MD5
        md5 = storage.calculate_md5(file_obj)
        file_obj.seek(0)
        
        # 检查是否已存在相同文件（MD5 + 大小双重校验）
        existing = await db.execute(
            select(cls.model).where(
                cls.model.md5 == md5,
                cls.model.size == file_size,
                cls.model.type == 'file',
                cls.model.is_deleted == False,  # noqa: E712
            ).limit(1)
        )
        existing_file = existing.scalar_one_or_none()
        
        if existing_file:
            # 复用已有文件的存储路径，不重复保存
            storage_path = existing_file.storage_path
            url = existing_file.url
            import logging
            logging.getLogger(__name__).info(f"文件 MD5 重复，复用已有存储路径: {storage_path}")
        else:
            # 新文件，保存到存储后端
            storage_path, url = storage.save(file_obj, filename, folder_path)
            import logging
            logging.getLogger(__name__).info(f"新文件保存到: {storage_path}")
        
        # 构建完整路径
        full_path = os.path.join(folder_path, filename).replace('\\', '/') if folder_path else filename
        
        # 创建数据库记录（即使文件内容相同，也创建独立记录，可能在不同文件夹、不同文件名）
        file_record = FileManager(
            name=filename,
            type='file',
            parent_id=parent_id,
            path=full_path,
            size=file_size,
            file_ext=file_ext,
            mime_type=mime_type,
            storage_type=storage.__class__.__name__.replace('StorageBackend', '').lower(),
            storage_path=storage_path,
            url=url,
            md5=md5,
            is_public=is_public,
            source=source,
            sys_creator_id=creator_id,
        )
        db.add(file_record)
        await db.commit()
        await db.refresh(file_record)
        return file_record

    @classmethod
    async def rename_item(
        cls,
        db: AsyncSession,
        item_id: str,
        new_name: str,
        creator_id: Optional[str] = None,
        is_superuser: bool = False,
    ) -> Optional[FileManager]:
        """重命名文件/文件夹
        
        权限规则：普通用户只能重命名自己创建的文件/文件夹
        """
        item = await cls.get_by_id(db, item_id)
        if not item:
            return None
        
        # 系统文件夹不允许重命名
        if item.is_system:
            return None
        
        # 权限校验：普通用户只能操作自己的文件
        if not is_superuser and creator_id and item.sys_creator_id != creator_id:
            return None
        
        # 检查同级目录下是否有同名文件
        existing = await db.execute(
            select(cls.model).where(
                cls.model.parent_id == item.parent_id,
                cls.model.name == new_name,
                cls.model.type == item.type,
                cls.model.id != item_id,
                cls.model.is_deleted == False  # noqa: E712
            )
        )
        if existing.scalar_one_or_none():
            return None  # 同名文件/文件夹已存在
        
        # 更新名称和路径
        old_path = item.path
        if item.parent_id:
            parent = await cls.get_by_id(db, item.parent_id)
            new_path = os.path.join(parent.path, new_name).replace('\\', '/') if parent else new_name
        else:
            new_path = new_name
        
        item.name = new_name
        item.path = new_path
        modifier_id = get_current_user_id_from_context()
        if modifier_id:
            item.sys_modifier_id = modifier_id
        await db.commit()
        await db.refresh(item)
        
        # 如果是文件夹，递归更新子项路径
        if item.type == 'folder':
            await cls._update_children_paths(db, item.id, old_path, new_path)
        
        return item

    @classmethod
    async def move_items(
        cls,
        db: AsyncSession,
        item_ids: List[str],
        target_folder_id: Optional[str] = None,
        creator_id: Optional[str] = None,
        is_superuser: bool = False,
    ) -> bool:
        """移动文件/文件夹
        
        权限规则：普通用户只能移动自己创建的文件/文件夹
        """
        # 获取目标文件夹
        target_path = ''
        if target_folder_id:
            target_folder = await cls.get_by_id(db, target_folder_id)
            if not target_folder or target_folder.type != 'folder':
                return False
            target_path = target_folder.path
        
        for item_id in item_ids:
            item = await cls.get_by_id(db, item_id)
            if not item:
                continue
            
            # 权限校验：普通用户只能操作自己的文件
            if not is_superuser and creator_id and item.sys_creator_id != creator_id:
                continue
            
            # 不能移动到自己或子文件夹
            if item.type == 'folder' and target_folder_id:
                if await cls._is_subfolder(db, target_folder_id, item.id):
                    continue
            
            # 检查目标文件夹是否有同名文件
            existing = await db.execute(
                select(cls.model).where(
                    cls.model.parent_id == target_folder_id,
                    cls.model.name == item.name,
                    cls.model.type == item.type,
                    cls.model.id != item_id,
                    cls.model.is_deleted == False  # noqa: E712
                )
            )
            if existing.scalar_one_or_none():
                continue
            
            # 更新父文件夹和路径
            old_path = item.path
            item.parent_id = target_folder_id
            item.path = os.path.join(target_path, item.name).replace('\\', '/') if target_path else item.name
            modifier_id = get_current_user_id_from_context()
            if modifier_id:
                item.sys_modifier_id = modifier_id
            
            # 如果是文件夹，递归更新子项路径
            if item.type == 'folder':
                await cls._update_children_paths(db, item.id, old_path, item.path)
        
        await db.commit()
        return True

    @classmethod
    async def delete_item(
        cls,
        db: AsyncSession,
        item_id: str,
        hard: bool = False,
        creator_id: Optional[str] = None,
        is_superuser: bool = False,
    ) -> bool:
        """删除文件/文件夹（默认软删除）
        
        权限规则：普通用户只能删除自己创建的文件/文件夹
        """
        item = await cls.get_by_id(db, item_id)
        if not item:
            return False
        
        # 系统文件夹不允许删除
        if item.is_system:
            return False
        
        # 权限校验：普通用户只能操作自己的文件
        if not is_superuser and creator_id and item.sys_creator_id != creator_id:
            return False
        
        # 如果是文件，检查是否有其他记录引用同一存储路径，没有才删除物理文件
        if item.type == 'file' and item.storage_path:
            ref_count = await db.execute(
                select(func.count(cls.model.id)).where(
                    cls.model.storage_path == item.storage_path,
                    cls.model.id != item.id,
                    cls.model.type == 'file',
                    cls.model.is_deleted == False,  # noqa: E712
                )
            )
            if (ref_count.scalar() or 0) == 0:
                storage = get_storage_backend()
                storage.delete(item.storage_path)
        
        # 递归删除子项
        if item.type == 'folder':
            await cls._delete_children(db, item.id, hard)
        
        # 删除数据库记录
        if hard:
            await db.delete(item)
        else:
            item.is_deleted = True
        
        await db.commit()
        return True

    @classmethod
    async def batch_delete(
        cls,
        db: AsyncSession,
        item_ids: List[str],
        hard: bool = False,
        creator_id: Optional[str] = None,
        is_superuser: bool = False,
    ) -> int:
        """批量删除文件/文件夹"""
        deleted_count = 0
        for item_id in item_ids:
            if await cls.delete_item(db, item_id, hard, creator_id=creator_id, is_superuser=is_superuser):
                deleted_count += 1
        return deleted_count

    @classmethod
    async def get_by_storage_path(
        cls,
        db: AsyncSession,
        storage_path: str,
    ) -> Optional[FileManager]:
        """通过存储路径获取文件"""
        result = await db.execute(
            select(cls.model).where(
                cls.model.storage_path == storage_path,
                cls.model.type == 'file',
                cls.model.is_deleted == False  # noqa: E712
            )
        )
        return result.scalar_one_or_none()

    @classmethod
    async def increment_download_count(
        cls,
        db: AsyncSession,
        item_id: str,
    ) -> None:
        """增加下载次数"""
        item = await cls.get_by_id(db, item_id)
        if item:
            item.download_count += 1
            await db.commit()

    @classmethod
    async def get_by_md5(
        cls,
        db: AsyncSession,
        md5: str,
        size: int,
    ) -> Optional[FileManager]:
        """通过MD5和大小查找文件（用于秒传）"""
        result = await db.execute(
            select(cls.model).where(
                cls.model.md5 == md5,
                cls.model.size == size,
                cls.model.is_deleted == False  # noqa: E712
            )
        )
        return result.scalar_one_or_none()

    @classmethod
    async def has_children(cls, db: AsyncSession, folder_id: str) -> bool:
        """检查文件夹是否有子项"""
        result = await db.execute(
            select(func.count(cls.model.id)).where(
                cls.model.parent_id == folder_id,
                cls.model.is_deleted == False  # noqa: E712
            )
        )
        count = result.scalar() or 0
        return count > 0

    @classmethod
    async def batch_has_children(cls, db: AsyncSession, folder_ids: List[str]) -> dict:
        """批量检查文件夹是否有子项，返回 {folder_id: bool}"""
        if not folder_ids:
            return {}
        result = await db.execute(
            select(
                cls.model.parent_id,
                func.count(cls.model.id).label('cnt')
            ).where(
                cls.model.parent_id.in_(folder_ids),
                cls.model.is_deleted == False  # noqa: E712
            ).group_by(cls.model.parent_id)
        )
        counts = {row[0]: row[1] > 0 for row in result.all()}
        return {fid: counts.get(fid, False) for fid in folder_ids}

    @classmethod
    async def batch_has_sub_folders(cls, db: AsyncSession, folder_ids: List[str]) -> dict:
        """批量检查文件夹是否有子文件夹（不含文件），返回 {folder_id: bool}"""
        if not folder_ids:
            return {}
        result = await db.execute(
            select(
                cls.model.parent_id,
                func.count(cls.model.id).label('cnt')
            ).where(
                cls.model.parent_id.in_(folder_ids),
                cls.model.type == 'folder',
                cls.model.is_deleted == False  # noqa: E712
            ).group_by(cls.model.parent_id)
        )
        counts = {row[0]: row[1] > 0 for row in result.all()}
        return {fid: counts.get(fid, False) for fid in folder_ids}

    @classmethod
    async def batch_get_names(cls, db: AsyncSession, item_ids: List[str]) -> dict:
        """批量获取文件/文件夹名称，返回 {id: name}"""
        if not item_ids:
            return {}
        result = await db.execute(
            select(cls.model.id, cls.model.name).where(
                cls.model.id.in_(item_ids),
                cls.model.is_deleted == False  # noqa: E712
            )
        )
        return {row[0]: row[1] for row in result.all()}

    @classmethod
    async def get_parent(cls, db: AsyncSession, item_id: str) -> Optional[FileManager]:
        """获取父文件夹"""
        item = await cls.get_by_id(db, item_id)
        if item and item.parent_id:
            return await cls.get_by_id(db, item.parent_id)
        return None

    @classmethod
    async def _is_subfolder(cls, db: AsyncSession, folder_id: str, potential_parent_id: str) -> bool:
        """检查folder是否是potential_parent的子文件夹（使用path前缀匹配）"""
        folder = await cls.get_by_id(db, folder_id)
        parent = await cls.get_by_id(db, potential_parent_id)
        if not folder or not parent:
            return False
        # 如果目标文件夹的path以潜在父文件夹的path为前缀，则是子文件夹
        parent_prefix = parent.path + '/'
        return folder.path.startswith(parent_prefix) or folder.id == potential_parent_id

    @classmethod
    async def _update_children_paths(cls, db: AsyncSession, folder_id: str, old_path: str, new_path: str) -> None:
        """批量更新所有后代路径（使用path前缀匹配，一条SQL搞定）"""
        old_prefix = old_path + '/'
        new_prefix = new_path + '/'
        # 使用 LIKE 前缀匹配找到所有后代，批量替换路径前缀
        stmt = (
            update(cls.model)
            .where(
                cls.model.path.like(f"{old_prefix}%"),
                cls.model.is_deleted == False  # noqa: E712
            )
            .values(
                path=func.concat(new_prefix, func.substr(cls.model.path, len(old_prefix) + 1))
            )
        )
        await db.execute(stmt)

    @classmethod
    async def _delete_children(cls, db: AsyncSession, folder_id: str, hard: bool = True) -> None:
        """批量删除所有后代（使用path前缀匹配查找所有后代）"""
        # 先获取当前文件夹的path
        folder = await cls.get_by_id(db, folder_id)
        if not folder:
            return
        
        folder_prefix = folder.path + '/'
        
        # 一次性查出所有后代文件（用于删除存储文件）
        result = await db.execute(
            select(cls.model).where(
                cls.model.path.like(f"{folder_prefix}%"),
                cls.model.type == 'file',
                cls.model.is_deleted == False  # noqa: E712
            )
        )
        file_children = result.scalars().all()
        
        # 删除存储文件
        if file_children:
            storage = get_storage_backend()
            for child in file_children:
                storage.delete(child.storage_path)
        
        # 批量更新/删除所有后代的数据库记录
        if hard:
            from sqlalchemy import delete as sql_delete
            stmt = sql_delete(cls.model).where(
                cls.model.path.like(f"{folder_prefix}%"),
                cls.model.is_deleted == False  # noqa: E712
            )
            await db.execute(stmt)
        else:
            stmt = (
                update(cls.model)
                .where(
                    cls.model.path.like(f"{folder_prefix}%"),
                    cls.model.is_deleted == False  # noqa: E712
                )
                .values(is_deleted=True)
            )
            await db.execute(stmt)

    # 默认最大允许全量读取的文件大小（200MB）
    MAX_FILE_CONTENT_SIZE = 200 * 1024 * 1024

    @classmethod
    async def get_file_content(cls, db: AsyncSession, file_id: str, max_size: int = None) -> Optional[bytes]:
        """
        获取文件内容（字节）
        
        Args:
            db: 数据库会话
            file_id: 文件ID
            max_size: 最大允许读取的文件大小（字节），超过则返回 None，默认200MB
            
        Returns:
            文件内容字节，如果文件不存在或超过大小限制则返回 None
        """
        file_obj = await cls.get_by_id(db, file_id)
        if not file_obj or file_obj.type != 'file':
            return None
        
        # 大文件保护
        limit = max_size or cls.MAX_FILE_CONTENT_SIZE
        if file_obj.size and file_obj.size > limit:
            import logging
            logging.warning(f"File {file_id} ({file_obj.name}) size {file_obj.size} exceeds max_size {limit}, skipping full read")
            return None
        
        storage = get_storage_backend()
        
        try:
            if isinstance(storage, LocalStorageBackend):
                # 本地存储：直接读取文件
                full_path = storage.get_full_path(file_obj.storage_path)
                if os.path.exists(full_path):
                    with open(full_path, 'rb') as f:
                        return f.read()
            elif isinstance(storage, MinioStorageBackend):
                # Minio 存储：通过 API 获取
                response = storage.get_file_content(file_obj.storage_path)
                content = response.read()
                response.close()
                response.release_conn()
                return content
            else:
                # 其他存储后端（OSS、Azure）：通过 URL 下载
                # 这里可以根据需要扩展
                import httpx
                url = file_obj.url
                if url:
                    async with httpx.AsyncClient() as client:
                        resp = await client.get(url)
                        if resp.status_code == 200:
                            return resp.content
        except Exception as e:
            import logging
            logging.error(f"Failed to get file content for {file_id}: {e}")
        
        return None

    @classmethod
    async def get_file_as_base64(cls, db: AsyncSession, file_id: str) -> Optional[str]:
        """
        获取文件内容的 base64 编码
        
        Args:
            db: 数据库会话
            file_id: 文件ID
            
        Returns:
            base64 编码的文件内容，如果文件不存在则返回 None
        """
        content = await cls.get_file_content(db, file_id)
        if content:
            return base64.b64encode(content).decode('utf-8')
        return None

    @classmethod
    async def get_file_as_data_url(cls, db: AsyncSession, file_id: str) -> Optional[str]:
        """
        获取文件的 Data URL（用于 LLM 多模态输入）
        
        格式: data:<mime_type>;base64,<base64_content>
        
        Args:
            db: 数据库会话
            file_id: 文件ID
            
        Returns:
            Data URL 格式的文件内容，如果文件不存在则返回 None
        """
        file_obj = await cls.get_by_id(db, file_id)
        if not file_obj or file_obj.type != 'file':
            return None
        
        content = await cls.get_file_content(db, file_id)
        if content:
            base64_content = base64.b64encode(content).decode('utf-8')
            mime_type = file_obj.mime_type or 'application/octet-stream'
            return f"data:{mime_type};base64,{base64_content}"
        return None

    # 支持提取文本内容的文件扩展名
    TEXT_EXTRACTABLE_EXTENSIONS = {
        # 纯文本
        '.txt', '.md', '.markdown', '.rst', '.log',
        # 代码文件
        '.py', '.js', '.ts', '.jsx', '.tsx', '.vue', '.html', '.css', '.scss', '.less',
        '.java', '.c', '.cpp', '.h', '.hpp', '.cs', '.go', '.rs', '.rb', '.php',
        '.swift', '.kt', '.scala', '.r', '.sql', '.sh', '.bash', '.zsh', '.ps1',
        '.yaml', '.yml', '.json', '.xml', '.toml', '.ini', '.cfg', '.conf',
        # 数据文件
        '.csv',
    }

    @classmethod
    async def get_file_text_content(
        cls, 
        db: AsyncSession, 
        file_id: str,
        max_size: int = 100 * 1024,  # 默认最大 100KB
        enable_ocr: bool = False,  # 是否启用 OCR（图片识别、扫描版 PDF）
    ) -> Optional[str]:
        """
        获取文件的文本内容（用于 LLM 处理）
        
        支持的文件类型：
        - 纯文本文件（txt, md, log 等）
        - 代码文件（py, js, ts, java 等）
        - 配置文件（json, yaml, xml 等）
        - CSV 文件
        - PDF 文件（需要 pypdf 库）
        - Word 文档（需要 python-docx 库）
        - 图片文件（需要 enable_ocr=True，通过视觉模型 OCR 识别）
        - 扫描版 PDF（需要 enable_ocr=True，自动 fallback 到 OCR）
        
        Args:
            db: 数据库会话
            file_id: 文件ID
            max_size: 最大读取大小（字节），超过则截断
            enable_ocr: 是否启用 OCR 识别（图片和扫描版 PDF）
            
        Returns:
            文件文本内容，如果无法提取则返回 None
        """
        file_obj = await cls.get_by_id(db, file_id)
        if not file_obj or file_obj.type != 'file':
            return None
        
        file_ext = (file_obj.file_ext or '').lower()
        mime_type = file_obj.mime_type or ''
        
        # 获取文件内容
        content = await cls.get_file_content(db, file_id)
        if not content:
            return None
        
        try:
            # 纯文本和代码文件
            if file_ext in cls.TEXT_EXTRACTABLE_EXTENSIONS or mime_type.startswith('text/'):
                text = content.decode('utf-8', errors='ignore')
                if len(text) > max_size:
                    text = text[:max_size] + f"\n\n... [内容已截断，原文件大小: {len(content)} 字节]"
                return text
            
            # PDF 文件
            if file_ext == '.pdf' or mime_type == 'application/pdf':
                try:
                    import io
                    from pypdf import PdfReader
                    reader = PdfReader(io.BytesIO(content))
                    text_parts = []
                    total_length = 0
                    for page in reader.pages:
                        page_text = page.extract_text() or ''
                        if total_length + len(page_text) > max_size:
                            text_parts.append(page_text[:max_size - total_length])
                            text_parts.append(f"\n\n... [内容已截断，共 {len(reader.pages)} 页]")
                            break
                        text_parts.append(page_text)
                        total_length += len(page_text)
                    
                    extracted_text = '\n'.join(text_parts).strip()
                    
                    # 检查是否为扫描版 PDF（提取的文本为空或几乎为空）
                    if len(extracted_text) < 10:
                        if enable_ocr:
                            # 扫描版 PDF：通过 OCR 识别
                            import logging
                            logging.info(f"PDF OCR: 扫描版 PDF {file_obj.name}，尝试 OCR 识别 ({len(reader.pages)} 页)")
                            ocr_text = await cls._ocr_pdf_pages(db, content, file_obj.name, max_size)
                            if ocr_text:
                                return ocr_text
                        return f"[PDF 文件: {file_obj.name}，共 {len(reader.pages)} 页]\n\n注意：该 PDF 文件可能是扫描版（图片格式），无法直接提取文本内容。"
                    
                    return extracted_text
                except ImportError:
                    # 返回提示信息，告知需要安装库
                    return f"[PDF 文件: {file_obj.name}]\n注意：需要安装 pypdf 库才能提取 PDF 内容。\n安装命令: pip install pypdf"
                except Exception as e:
                    # 返回错误信息
                    import logging
                    logging.error(f"Failed to extract PDF content from {file_obj.name}: {e}")
                    return f"[PDF 文件: {file_obj.name}]\n提取内容时出错: {str(e)}"
            
            # Word 文档
            if file_ext in ('.docx', '.doc') or mime_type in (
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                'application/msword'
            ):
                try:
                    import io
                    from docx import Document
                    doc = Document(io.BytesIO(content))
                    text_parts = []
                    total_length = 0
                    for para in doc.paragraphs:
                        para_text = para.text
                        if total_length + len(para_text) > max_size:
                            text_parts.append(para_text[:max_size - total_length])
                            text_parts.append("\n\n... [内容已截断]")
                            break
                        text_parts.append(para_text)
                        total_length += len(para_text)
                    return '\n'.join(text_parts)
                except ImportError:
                    return f"[Word 文档: {file_obj.name}，需要安装 python-docx 库才能提取内容]"
                except Exception as e:
                    return f"[Word 文档: {file_obj.name}，提取内容失败: {str(e)}]"
            
            # Excel 文件
            if file_ext in ('.xlsx', '.xls') or mime_type in (
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                'application/vnd.ms-excel'
            ):
                try:
                    import io
                    from openpyxl import load_workbook
                    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
                    ws = wb.active
                    text_parts = []
                    total_length = 0
                    for row in ws.iter_rows(values_only=True):
                        row_text = ','.join([str(cell) if cell is not None else '' for cell in row])
                        if total_length + len(row_text) > max_size:
                            text_parts.append(row_text[:max_size - total_length])
                            text_parts.append(f"\n\n... [内容已截断，原表格共 {ws.max_row} 行]")
                            break
                        text_parts.append(row_text)
                        total_length += len(row_text)
                    wb.close()
                    return '\n'.join(text_parts)
                except ImportError:
                    return f"[Excel 文件: {file_obj.name}]\n注意：需要安装 openpyxl 库才能提取 Excel 内容。\n安装命令: pip install openpyxl"
                except Exception as e:
                    import logging
                    logging.error(f"Failed to extract Excel content from {file_obj.name}: {e}")
                    return f"[Excel 文件: {file_obj.name}]\n提取内容时出错: {str(e)}"
            
            # PowerPoint 文件
            if file_ext in ('.pptx', '.ppt') or mime_type in (
                'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'application/vnd.ms-powerpoint'
            ):
                try:
                    import io
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(content))
                    text_parts = []
                    total_length = 0
                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_text = f"--- 幻灯片 {slide_num} ---\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slide_text += shape.text + "\n"
                        if total_length + len(slide_text) > max_size:
                            text_parts.append(slide_text[:max_size - total_length])
                            text_parts.append(f"\n\n... [内容已截断，共 {len(prs.slides)} 页幻灯片]")
                            break
                        text_parts.append(slide_text)
                        total_length += len(slide_text)
                    return '\n'.join(text_parts)
                except ImportError:
                    return f"[PowerPoint 文件: {file_obj.name}]\n注意：需要安装 python-pptx 库才能提取 PPT 内容。\n安装命令: pip install python-pptx"
                except Exception as e:
                    import logging
                    logging.error(f"Failed to extract PowerPoint content from {file_obj.name}: {e}")
                    return f"[PowerPoint 文件: {file_obj.name}]\n提取内容时出错: {str(e)}"
            
            # 图片文件：通过 OCR 识别
            if enable_ocr and (file_ext in cls.IMAGE_EXTENSIONS or mime_type.startswith('image/')):
                import logging
                logging.info(f"Image OCR: 识别图片 {file_obj.name}")
                ocr_text = await cls._ocr_from_image_bytes(db, content, mime_type or 'image/png')
                return ocr_text
            
            # 不支持的文件类型
            return None
            
        except Exception as e:
            import logging
            logging.error(f"Failed to extract text from file {file_id}: {e}")
            return None

    @classmethod
    async def recognize_image_with_ocr(
        cls,
        db: AsyncSession,
        file_id: str,
        prompt: str = "请识别并提取图片中的所有文字内容，保持原有格式和结构。",
    ) -> Optional[str]:
        """
        使用阿里云 qwen-vl-ocr 模型识别图片中的文字内容
        
        用于非多模态模型处理图片附件时，先通过 OCR 提取图片中的文字
        
        Args:
            db: 数据库会话
            file_id: 文件ID
            prompt: OCR 识别的提示词
            
        Returns:
            识别出的文字内容，如果识别失败则返回 None
        """
        import logging
        import os
        
        # 获取文件信息
        file_obj = await cls.get_by_id(db, file_id)
        if not file_obj or file_obj.type != 'file':
            logging.warning(f"OCR: File not found or not a file: {file_id}")
            return None
        
        # 检查是否为图片
        mime_type = file_obj.mime_type or ''
        if not mime_type.startswith('image/'):
            logging.warning(f"OCR: File is not an image: {file_obj.name}, mime_type={mime_type}")
            return None
        
        # 获取文件内容并转换为 base64
        content = await cls.get_file_content(db, file_id)
        if not content:
            logging.warning(f"OCR: Failed to get file content for file: {file_id}")
            return None
        
        # 转换为 base64 字符串（qwen-vl-ocr 需要 data URL 格式）
        import base64
        base64_str = base64.b64encode(content).decode('utf-8')
        mime_type = file_obj.mime_type or 'image/png'
        base64_data_url = f"data:{mime_type};base64,{base64_str}"
        
        try:
            from openai import AsyncOpenAI
            
            # 通过统一辅助方法获取 API 配置（数据库 qwen provider 优先，env 兜底）
            ocr_config = await cls._get_ocr_client(db)
            if not ocr_config:
                return f"[图片: {file_obj.name}]\n注意：未配置阿里云 API Key，无法进行 OCR 识别。"
            
            api_key, api_base, _ = ocr_config
            
            client = AsyncOpenAI(
                api_key=api_key,
                base_url=api_base,
            )
            
            logging.info(f"OCR: Recognizing image {file_obj.name} with qwen-vl-ocr")
            
            response = await client.chat.completions.create(
                model="qwen-vl-ocr-latest",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "image_url",
                                "image_url": {"url": base64_data_url},
                            },
                            {"type": "text", "text": prompt}
                        ]
                    }
                ]
            )
            
            if response.choices and response.choices[0].message.content:
                ocr_text = response.choices[0].message.content
                logging.info(f"OCR: Successfully recognized {len(ocr_text)} characters from {file_obj.name}")
                return ocr_text
            else:
                logging.warning(f"OCR: No content in response for {file_obj.name}")
                return None
                
        except ImportError:
            logging.error("OCR: openai library not installed")
            return f"[图片: {file_obj.name}]\n注意：需要安装 openai 库才能进行 OCR 识别。"
        except Exception as e:
            logging.error(f"OCR: Failed to recognize image {file_obj.name}: {e}")
            return f"[图片: {file_obj.name}]\nOCR 识别失败: {str(e)}"

    @classmethod
    async def _get_ocr_client(cls, db: AsyncSession):
        """
        获取 OCR 所需的 API 客户端配置
        
        优先级：
        1. 数据库中 qwen 类型提供商的 API Key
        2. 环境变量 / settings 中的 DASHSCOPE_API_KEY（兜底）
        
        Returns:
            (api_key, api_base, model_name) 或 None
        """
        import logging
        from app.config import settings

        api_key = None
        api_base = "https://dashscope.aliyuncs.com/compatible-mode/v1"
        model_name = "qwen-vl-ocr-latest"

        # 1. 优先从数据库获取 qwen 提供商的 API Key
        try:
            from sqlalchemy import select as sa_select
            from ai_platform.models.provider import LLMProvider
            provider_result = await db.execute(
                sa_select(LLMProvider).where(
                    LLMProvider.provider_type == 'qwen',
                    LLMProvider.is_active == True,
                    LLMProvider.is_deleted == False
                )
            )
            provider = provider_result.scalar_one_or_none()
            if provider and provider.api_key:
                api_key = provider.api_key
                if provider.api_base:
                    api_base = provider.api_base
                logging.info(f"OCR: Using API key from database provider: {provider.name}")
        except Exception as e:
            logging.debug(f"OCR: Failed to get provider from database: {e}")

        # 2. 兜底：从 settings 获取（settings 已自动从环境变量读取）
        if not api_key:
            api_key = getattr(settings, 'DASHSCOPE_API_KEY', None)

        if not api_key:
            logging.warning("OCR: No API key available for OCR")
            return None

        return api_key, api_base, model_name

    @classmethod
    async def _ocr_from_image_bytes(
        cls,
        db: AsyncSession,
        image_bytes: bytes,
        mime_type: str = 'image/png',
        prompt: str = "请识别并提取图片中的所有文字内容，保持原有格式和结构。",
    ) -> Optional[str]:
        """
        从图片字节数据中 OCR 提取文字
        
        Args:
            db: 数据库会话
            image_bytes: 图片文件字节
            mime_type: 图片 MIME 类型
            prompt: OCR 提示词
            
        Returns:
            识别出的文字内容
        """
        import logging

        config = await cls._get_ocr_client(db)
        if not config:
            return None

        api_key, api_base, model_name = config

        try:
            from openai import AsyncOpenAI

            base64_str = base64.b64encode(image_bytes).decode('utf-8')
            data_url = f"data:{mime_type};base64,{base64_str}"

            client = AsyncOpenAI(api_key=api_key, base_url=api_base)
            response = await client.chat.completions.create(
                model=model_name,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image_url", "image_url": {"url": data_url}},
                        {"type": "text", "text": prompt},
                    ]
                }]
            )

            if response.choices and response.choices[0].message.content:
                text = response.choices[0].message.content
                logging.info(f"OCR: 识别到 {len(text)} 个字符")
                return text
            return None
        except ImportError:
            logging.error("OCR: openai library not installed")
            return None
        except Exception as e:
            logging.error(f"OCR: 图片识别失败: {e}")
            return None

    @classmethod
    async def _ocr_pdf_pages(
        cls,
        db: AsyncSession,
        pdf_bytes: bytes,
        filename: str,
        max_size: int = 100 * 1024,
        max_pages: int = 20,
    ) -> Optional[str]:
        """
        对扫描版 PDF 逐页进行 OCR 识别
        
        尝试两种策略：
        1. 使用 pypdf 提取页面中嵌入的图片
        2. 使用 pdf2image 将页面渲染为图片（需要 poppler）
        
        Args:
            db: 数据库会话
            pdf_bytes: PDF 文件字节
            filename: 文件名（用于日志）
            max_size: 最大文本大小
            max_pages: 最大处理页数
            
        Returns:
            OCR 识别的文本内容
        """
        import io
        import logging

        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(pdf_bytes))
            total_pages = len(reader.pages)
            pages_to_process = min(total_pages, max_pages)

            text_parts = []
            total_length = 0

            for page_idx in range(pages_to_process):
                page = reader.pages[page_idx]

                # 尝试从页面提取嵌入图片
                page_image = cls._extract_largest_image_from_page(page)
                if not page_image:
                    continue

                img_bytes, img_mime = page_image
                ocr_text = await cls._ocr_from_image_bytes(
                    db, img_bytes, img_mime,
                    prompt=f"请识别并提取这张图片（PDF第{page_idx + 1}页）中的所有文字内容，保持原有格式和结构。",
                )

                if ocr_text:
                    page_header = f"--- 第 {page_idx + 1} 页 ---\n"
                    page_content = page_header + ocr_text

                    if total_length + len(page_content) > max_size:
                        text_parts.append(page_content[:max_size - total_length])
                        text_parts.append(f"\n\n... [内容已截断，共 {total_pages} 页，已处理 {page_idx + 1} 页]")
                        break
                    text_parts.append(page_content)
                    total_length += len(page_content)

            if text_parts:
                result = '\n\n'.join(text_parts)
                logging.info(f"PDF OCR: {filename} 识别完成，{len(text_parts)} 页，{len(result)} 字符")
                return result

            # 如果 pypdf 无法提取图片，尝试 pdf2image
            return await cls._ocr_pdf_with_pdf2image(db, pdf_bytes, filename, max_size, max_pages)

        except Exception as e:
            logging.error(f"PDF OCR: {filename} 处理失败: {e}")
            return None

    @classmethod
    def _extract_largest_image_from_page(cls, page) -> Optional[tuple]:
        """
        从 PDF 页面中提取最大的嵌入图片
        
        Returns:
            (image_bytes, mime_type) 或 None
        """
        try:
            if not hasattr(page, 'images') or not page.images:
                return None

            largest = None
            largest_size = 0
            for img in page.images:
                img_data = img.data
                if len(img_data) > largest_size:
                    largest_size = len(img_data)
                    # 根据图片名称推断 MIME 类型
                    name = (img.name or '').lower()
                    if name.endswith('.png'):
                        mime = 'image/png'
                    elif name.endswith('.jpg') or name.endswith('.jpeg'):
                        mime = 'image/jpeg'
                    else:
                        mime = 'image/png'
                    largest = (img_data, mime)

            return largest
        except Exception:
            return None

    @classmethod
    async def _ocr_pdf_with_pdf2image(
        cls,
        db: AsyncSession,
        pdf_bytes: bytes,
        filename: str,
        max_size: int = 100 * 1024,
        max_pages: int = 20,
    ) -> Optional[str]:
        """
        使用 pdf2image 将 PDF 页面渲染为图片后 OCR（需要 poppler）
        """
        import io
        import logging

        try:
            from pdf2image import convert_from_bytes
        except ImportError:
            logging.warning("PDF OCR: pdf2image 未安装，无法渲染扫描版 PDF。安装命令: pip install pdf2image")
            return None

        try:
            images = convert_from_bytes(
                pdf_bytes,
                first_page=1,
                last_page=max_pages,
                dpi=200,
            )

            text_parts = []
            total_length = 0

            for idx, img in enumerate(images):
                # 将 PIL Image 转为 PNG 字节
                buf = io.BytesIO()
                img.save(buf, format='PNG')
                img_bytes = buf.getvalue()

                ocr_text = await cls._ocr_from_image_bytes(
                    db, img_bytes, 'image/png',
                    prompt=f"请识别并提取这张图片（PDF第{idx + 1}页）中的所有文字内容，保持原有格式和结构。",
                )

                if ocr_text:
                    page_header = f"--- 第 {idx + 1} 页 ---\n"
                    page_content = page_header + ocr_text

                    if total_length + len(page_content) > max_size:
                        text_parts.append(page_content[:max_size - total_length])
                        text_parts.append(f"\n\n... [内容已截断，共 {len(images)} 页渲染]")
                        break
                    text_parts.append(page_content)
                    total_length += len(page_content)

            if text_parts:
                result = '\n\n'.join(text_parts)
                logging.info(f"PDF OCR (pdf2image): {filename} 识别完成，{len(text_parts)} 页，{len(result)} 字符")
                return result

            return None
        except Exception as e:
            logging.error(f"PDF OCR (pdf2image): {filename} 渲染失败: {e}")
            return None

    @classmethod
    def _build_function_config_from_schema(cls, output_schema: list) -> dict:
        """
        将前端传递的 output_schema 转换为 Function Calling 配置
        
        Args:
            output_schema: 前端传递的结构化输出字段定义列表
            
        Returns:
            Function Calling 配置字典
        """
        def build_properties(fields: list) -> tuple:
            """递归构建 properties 和 required 列表"""
            properties = {}
            required = []
            
            for field in fields:
                if not field.get("name"):
                    continue
                    
                prop = {
                    "type": field.get("type", "string"),
                    "description": field.get("description", ""),
                }
                
                # 处理枚举值
                if field.get("enum"):
                    prop["enum"] = field["enum"]
                
                # 处理对象类型
                if field.get("type") == "object" and field.get("properties"):
                    nested_props, nested_required = build_properties(field["properties"])
                    prop["properties"] = nested_props
                    if nested_required:
                        prop["required"] = nested_required
                
                # 处理数组类型
                elif field.get("type") == "array" and field.get("items"):
                    items = field["items"]
                    if items.get("type") == "object" and items.get("properties"):
                        nested_props, nested_required = build_properties(items["properties"])
                        prop["items"] = {
                            "type": "object",
                            "properties": nested_props,
                        }
                        if nested_required:
                            prop["items"]["required"] = nested_required
                    else:
                        prop["items"] = {"type": items.get("type", "string")}
                
                properties[field["name"]] = prop
                if field.get("required"):
                    required.append(field["name"])
            
            return properties, required
        
        properties, required = build_properties(output_schema)
        
        return {
            "name": "extract_structured_data",
            "description": "从文件内容中提取结构化数据",
            "parameters": {
                "type": "object",
                "properties": properties,
                "required": required if required else None,
            }
        }

    @classmethod
    async def recognize_file_with_function_calling(
        cls,
        db: AsyncSession,
        file_id: str,
        output_schema: list = None,
        custom_prompt: str = None,
    ) -> dict:
        """
        使用 AI 识别文件内容并通过 Function Calling 提取结构化数据
        
        支持的文件类型：
        - 图片文件：使用 qwen-vl-ocr 模型进行 OCR 识别
        - 文本文件：直接提取文本内容
        - PDF 文件：提取 PDF 文本内容
        - Word 文档：提取 Word 文本内容
        - Excel 文件：提取 Excel 内容
        
        Args:
            db: 数据库会话
            file_id: 文件ID
            output_schema: 结构化输出字段定义（前端传递）
            custom_prompt: 自定义提示词
            
        Returns:
            {
                "success": bool,
                "raw_text": str,  # 原始识别/提取的文字
                "extracted_data": dict,  # 提取的结构化数据
                "error": str  # 错误信息
            }
        """
        import logging
        import os
        import json
        
        # 如果有 output_schema，将其转换为 function_config
        function_config = None
        if output_schema:
            function_config = cls._build_function_config_from_schema(output_schema)
        
        # 获取文件信息
        file_obj = await cls.get_by_id(db, file_id)
        if not file_obj or file_obj.type != 'file':
            return {"success": False, "error": "文件不存在", "raw_text": None, "extracted_data": None}
        
        # 检查文件类型
        mime_type = file_obj.mime_type or ''
        file_ext = (file_obj.file_ext or '').lower()
        is_image = mime_type.startswith('image/')
        
        # 非图片文件：使用文本提取方法
        if not is_image:
            return await cls._recognize_text_file_with_function_calling(
                db, file_id, file_obj, function_config, custom_prompt
            )
        
        # 获取文件内容并转换为 base64
        content = await cls.get_file_content(db, file_id)
        if not content:
            return {"success": False, "error": "无法读取文件内容", "raw_text": None, "extracted_data": None}
        
        import base64
        base64_str = base64.b64encode(content).decode('utf-8')
        mime_type = file_obj.mime_type or 'image/png'
        base64_data_url = f"data:{mime_type};base64,{base64_str}"
        
        try:
            from openai import AsyncOpenAI
            
            # 通过统一辅助方法获取 API 配置（数据库 qwen provider 优先，env 兜底）
            ocr_config = await cls._get_ocr_client(db)
            if not ocr_config:
                return {"success": False, "error": "未配置阿里云 API Key", "raw_text": None, "extracted_data": None}
            
            api_key, api_base, _ = ocr_config
            client = AsyncOpenAI(api_key=api_key, base_url=api_base)
            
            # 使用自定义提示词或默认提示词
            prompt = custom_prompt or "请识别这个文件中的所有内容。"
            
            # 如果没有 function 配置，只做普通 OCR
            if not function_config:
                response = await client.chat.completions.create(
                    model="qwen-vl-ocr-latest",
                    messages=[{
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": {"url": base64_data_url}},
                            {"type": "text", "text": prompt}
                        ]
                    }]
                )
                
                if response.choices and response.choices[0].message.content:
                    return {
                        "success": True,
                        "raw_text": response.choices[0].message.content,
                        "extracted_data": None,
                        "error": None
                    }
                return {"success": False, "error": "识别结果为空", "raw_text": None, "extracted_data": None}
            
            # 使用 Function Calling
            tools = [{
                "type": "function",
                "function": {
                    "name": function_config["name"],
                    "description": function_config["description"],
                    "parameters": function_config["parameters"]
                }
            }]
            
            logging.info(f"OCR with Function Calling: {file_obj.name}")
            
            response = await client.chat.completions.create(
                model="qwen-vl-ocr-latest",
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image_url", "image_url": {"url": base64_data_url}},
                        {"type": "text", "text": f"{prompt}\n请调用 {function_config['name']} 函数提取结构化信息。"}
                    ]
                }],
                tools=tools,
                tool_choice="auto"
            )
            
            if not response.choices:
                return {"success": False, "error": "API 返回为空", "raw_text": None, "extracted_data": None}
            
            message = response.choices[0].message
            raw_text = message.content or ""
            extracted_data = None
            
            # 检查是否有 tool_calls
            if message.tool_calls:
                for tool_call in message.tool_calls:
                    if tool_call.function.name == function_config["name"]:
                        try:
                            extracted_data = json.loads(tool_call.function.arguments)
                            logging.info(f"OCR extracted data: {extracted_data}")
                        except json.JSONDecodeError as e:
                            logging.error(f"Failed to parse function arguments: {e}")
                        break
            
            return {
                "success": True,
                "raw_text": raw_text,
                "extracted_data": extracted_data,
                "error": None
            }
            
        except ImportError:
            return {"success": False, "error": "需要安装 openai 库", "raw_text": None, "extracted_data": None}
        except Exception as e:
            logging.error(f"OCR with Function Calling failed: {e}")
            return {"success": False, "error": str(e), "raw_text": None, "extracted_data": None}

    @classmethod
    async def _recognize_text_file_with_function_calling(
        cls,
        db: AsyncSession,
        file_id: str,
        file_obj,
        function_config: dict,
        custom_prompt: str,
    ) -> dict:
        """
        处理非图片文件的识别（文本、PDF、Word等）
        使用 LLM 进行内容理解和结构化提取
        """
        import logging
        import os
        import json
        
        # 提取文件文本内容
        text_content = await cls.get_file_text_content(db, file_id)
        if not text_content:
            return {"success": False, "error": "无法提取文件内容", "raw_text": None, "extracted_data": None}
        
        # 如果没有 function_config，直接返回文本内容
        if not function_config:
            return {
                "success": True,
                "raw_text": text_content,
                "extracted_data": None,
                "error": None
            }
        
        # 使用 LLM 进行结构化提取
        try:
            from openai import AsyncOpenAI
            from app.config import settings
            
            api_key = None
            api_base = None
            model = "gpt-4o-mini"
            
            # 1. 优先从数据库获取任意激活的 LLM 提供商
            try:
                from sqlalchemy import select
                from ai_platform.models.provider import LLMProvider
                provider_result = await db.execute(
                    select(LLMProvider).where(
                        LLMProvider.is_active == True,
                        LLMProvider.is_deleted == False
                    ).order_by(LLMProvider.sort.asc())
                )
                provider = provider_result.scalar_one_or_none()
                if provider and provider.api_key:
                    api_key = provider.api_key
                    api_base = provider.api_base
                    if provider.provider_type == 'qwen':
                        model = "qwen-plus"
                        if not api_base:
                            api_base = "https://dashscope.aliyuncs.com/compatible-mode/v1"
                    elif provider.provider_type == 'openai':
                        model = "gpt-4o-mini"
                    else:
                        model = "gpt-4o-mini"
            except Exception:
                pass
            
            # 2. 兜底：从 settings/环境变量获取
            if not api_key:
                api_key = getattr(settings, 'OPENAI_API_KEY', None) or os.getenv("OPENAI_API_KEY")
                api_base = getattr(settings, 'OPENAI_API_BASE', None) or os.getenv("OPENAI_API_BASE")
                model = getattr(settings, 'OPENAI_MODEL', None) or os.getenv("OPENAI_MODEL") or "gpt-4o-mini"
            
            if not api_key:
                api_key = getattr(settings, 'DASHSCOPE_API_KEY', None)
                if api_key:
                    api_base = "https://dashscope.aliyuncs.com/compatible-mode/v1"
                    model = "qwen-plus"
            
            if not api_key:
                return {"success": False, "error": "未配置 LLM API Key", "raw_text": text_content, "extracted_data": None}
            
            client = AsyncOpenAI(api_key=api_key, base_url=api_base)
            
            # 构建提示词
            prompt = custom_prompt or "请从以下文件内容中提取结构化信息。"
            
            tools = [{
                "type": "function",
                "function": {
                    "name": function_config["name"],
                    "description": function_config["description"],
                    "parameters": function_config["parameters"]
                }
            }]
            
            logging.info(f"Text file extraction with Function Calling: {file_obj.name}")
            
            response = await client.chat.completions.create(
                model=model,
                messages=[
                    {
                        "role": "system",
                        "content": "你是一个专业的文档信息提取助手。请仔细阅读用户提供的文件内容，并使用提供的函数提取结构化信息。"
                    },
                    {
                        "role": "user",
                        "content": f"{prompt}\n\n文件内容：\n{text_content}"
                    }
                ],
                tools=tools,
                tool_choice="auto"
            )
            
            if not response.choices:
                return {"success": False, "error": "API 返回为空", "raw_text": text_content, "extracted_data": None}
            
            message = response.choices[0].message
            extracted_data = None
            
            # 检查是否有 tool_calls
            if message.tool_calls:
                for tool_call in message.tool_calls:
                    if tool_call.function.name == function_config["name"]:
                        try:
                            extracted_data = json.loads(tool_call.function.arguments)
                            logging.info(f"Text file extracted data: {extracted_data}")
                        except json.JSONDecodeError as e:
                            logging.error(f"Failed to parse function arguments: {e}")
                        break
            
            return {
                "success": True,
                "raw_text": text_content,
                "extracted_data": extracted_data,
                "error": None
            }
            
        except ImportError:
            return {"success": False, "error": "需要安装 openai 库", "raw_text": text_content, "extracted_data": None}
        except Exception as e:
            logging.error(f"Text file extraction with Function Calling failed: {e}")
            return {"success": False, "error": str(e), "raw_text": text_content, "extracted_data": None}
    
    # 保留旧方法名作为别名，保持向后兼容
    recognize_image_with_function_calling = recognize_file_with_function_calling
